from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.util import Inches

import os.path as path
import os
import random

from moviepy.editor import VideoFileClip

import struct
import imghdr

def get_image_size(fname):
    '''Determine the image type of fhandle and return its size.
    from draco'''
    with open(fname, 'rb') as fhandle:
        head = fhandle.read(24)
        if len(head) != 24:
            return
        if imghdr.what(fname) == 'png':
            check = struct.unpack('>i', head[4:8])[0]
            if check != 0x0d0a1a0a:
                return
            width, height = struct.unpack('>ii', head[16:24])
        elif imghdr.what(fname) == 'gif':
            width, height = struct.unpack('<HH', head[6:10])
        elif imghdr.what(fname) == 'jpeg':
            try:
                fhandle.seek(0) # Read 0xff next
                size = 2
                ftype = 0
                while not 0xc0 <= ftype <= 0xcf:
                    fhandle.seek(size, 1)
                    byte = fhandle.read(1)
                    while ord(byte) == 0xff:
                        byte = fhandle.read(1)
                    ftype = ord(byte)
                    size = struct.unpack('>H', fhandle.read(2))[0] - 2
                # We are at a SOFn block
                fhandle.seek(1, 1)  # Skip `precision' byte.
                height, width = struct.unpack('>HH', fhandle.read(4))
            except Exception: #IGNORE:W0703
                return
        else:
            return
        return width, height

PPI=72


pathRuning=True
while pathRuning:
	print("Bienvenue sur PhotoPower")
	pathFile=path.normcase(path.normpath(input("Entrez le chemin du dossier à utiliser: ")))
	if path.exists(pathFile):
		fileListTemp=os.listdir(pathFile)
		fileList=[]
		for i in fileListTemp:
			if(imghdr.what("%s/%s"%(pathFile,i)) == 'png' or imghdr.what("%s/%s"%(pathFile,i)) == 'jpeg' or imghdr.what("%s/%s"%(pathFile,i)) == 'mp4'):
				fileList.append(i)
		if(len(fileList)==0):
			print("%s Ne contient pas d'éléments."%(pathFile))
		else:
			print("Le chemin %s a été sélectioné, il contient %d photo et vidéo."%(pathFile, len(fileList)))
			pathRuning=False
	else:
		print("Chemin incorect")

for i in fileList:
	fileList[fileList.index(i)]="%s/%s"%(pathFile, i)

nameRuning=True
while nameRuning:
	name=input("Veuillez saisir le nom de votre présentation: ")
	if(len(name)<3):
		print("Le nom doit contenir au moins 3 charactères")
	else:
		print("Le nom %s a été saisi"%(name))
		nameRuning=False

if len(fileList)%100==0:
	rangeNbr=len(fileList)/100
else:
	rangeNbr=(int)(len(fileList)/100+1)

for i in range(rangeNbr):
	prs = Presentation()

	titleLayout = prs.slide_layouts[0]
	blankLayout = prs.slide_layouts[6]

	slHeight = prs.slide_height
	slWidth = prs.slide_width
	ratio=slHeight/slWidth

	titleSlide = prs.slides.add_slide(titleLayout)
	titleSlide.shapes.title.text=name
	if(len(fileList)<100):
		for j in range(len(fileList)):
			titleSlide.placeholders[1].text="%d-%d"%(i*100+1,i*100+len(fileList))
			picturePath=random.choice(fileList)
			sl=prs.slides.add_slide(blankLayout)
			shape=sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,slWidth,slHeight)
			shape.fill.solid()
			shape.fill.fore_color.rgb=RGBColor(0,0,0)

			if(imghdr.what(picturePath) == 'png' or imghdr.what(picturePath) == 'jpeg'):
				pic=sl.shapes.add_picture(picturePath, 0, 0)

				w,h=get_image_size(picturePath)
				wDpi,hDpi=pic.image.dpi

				w=Inches(w/wDpi)
				h=Inches(h/hDpi)
				ratioImg=h/w
				if(ratio<ratioImg):
					pic.left=int(slWidth/2-(slHeight/ratioImg)/2)
					pic.height=slHeight
					pic.width=int(slHeight/ratioImg)
				else:
					pic.top=int(slHeight/2-(slWidth*ratioImg)/2)
					pic.width=slWidth
					pic.height=int(slWidth*ratioImg)
			elif(imghdr.what(picturePath) == 'mp4'):
				clip=VideoFileClip(picturePath)
				w,h=clip.size
				w=Inches(w/72)
				h=Inches(h/72)
				ratioImg=h/w

				clip.save_frame("./temp/poster-%d-%d.jpeg"%(i,j))

				if(ratio<ratioImg):
					sl.shapes.add_movie(picturePath, int(slWidth/2-(slHeight/ratioImg)/2),0,int(slHeight/ratioImg), slHeight,"./temp/poster.jpeg")
				else:
					sl.shapes.add_movie(picturePath, 0, int(slHeight/2-(slWidth*ratioImg)/2), slWidth, int(slWidth*ratioImg), "./temp/poster.jpeg")
			else:
				print(picturePath)

			fileList.remove(picturePath)
	else:
		for j in range(100):
			titleSlide.placeholders[1].text="%d-%d"%(i*100+1,(i+1)*100)
			picturePath=random.choice(fileList)
			sl=prs.slides.add_slide(blankLayout)
			shape=sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,slWidth,slHeight)
			shape.fill.solid()
			shape.fill.fore_color.rgb=RGBColor(0,0,0)

			if(imghdr.what(picturePath) == 'png' or imghdr.what(picturePath) == 'jpeg'):
				pic=sl.shapes.add_picture(picturePath, 0, 0)

				w,h=get_image_size(picturePath)
				wDpi,hDpi=pic.image.dpi

				w=Inches(w/wDpi)
				h=Inches(h/hDpi)
				ratioImg=h/w
				if(ratio<ratioImg):
					pic.left=int(slWidth/2-(slHeight/ratioImg)/2)
					pic.height=slHeight
					pic.width=int(slHeight/ratioImg)
				else:
					pic.top=int(slHeight/2-(slWidth*ratioImg)/2)
					pic.width=slWidth
					pic.height=int(slWidth*ratioImg)
			elif(imghdr.what(picturePath) == 'mp4'):
				clip=VideoFileClip(picturePath)
				w,h=clip.size
				w=Inches(w/72)
				h=Inches(h/72)
				ratioImg=h/w

				clip.save_frame("./temp/poster-%d-%d.jpeg"%(i,j))

				if(ratio<ratioImg):
					sl.shapes.add_movie(picturePath, int(slWidth/2-(slHeight/ratioImg)/2),0,int(slHeight/ratioImg), slHeight,"./temp/poster-%d-%d.jpeg"%(i,j))
				else:
					sl.shapes.add_movie(picturePath, 0, int(slHeight/2-(slWidth*ratioImg)/2), slWidth, int(slWidth*ratioImg), "./temp/poster-%d-%d.jpeg"%(i,j))
			else:
				print(picturePath)
			fileList.remove(picturePath)
	prs.save("%s/%s - %d.pptx"%(pathFile, name,i+1))
	print("Document %s/%s - %d.pptx créé"%(pathFile, name, i+1))

if(path.exists("./temp")):
	for i in os.listdir("./temp"):
		os.remove("./temp/%s"%(i))
